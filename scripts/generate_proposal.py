#!/usr/bin/env python3
"""proposal_data.json → PPTX生成スクリプト

使い方:
    python scripts/generate_proposal.py <proposal_data.json>

proposal_data.json の構造:
{
    "template": "shared/templates/pptx/tomoshi-base.pptx",  (省略時はブランクから生成)
    "output": "projects/xxx/proposals/generated/proposal_v1.pptx",
    "slides": [
        {"layout": "title", "title": "...", "subtitle": "..."},
        {"layout": "section", "title": "..."},
        {"layout": "content", "title": "...", "body": "- 項目1\n- **太字**項目2"},
        {"layout": "two_column", "title": "...", "left": "...", "right": "..."},
        {"layout": "table", "title": "...", "headers": [...], "rows": [[...], ...]},
        {"layout": "closing", "title": "...", "body": "..."}
    ]
}
"""

import json
import sys
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# scriptsディレクトリからインポート
sys.path.insert(0, str(Path(__file__).parent))
from utils import (
    parse_markdown_to_paragraphs,
    apply_paragraphs_to_textframe,
    set_text_with_font,
    FONT_NAME,
    FONT_SIZE_BODY,
    FONT_SIZE_BULLET,
    COLOR_PRIMARY,
    COLOR_ACCENT,
    COLOR_LIGHT,
)

# プロジェクトルート
PROJECT_ROOT = Path(__file__).parent.parent


def find_layout(prs, layout_name_candidates):
    """テンプレートからレイアウトを名前で探す。見つからなければインデックスで返す。"""
    for layout in prs.slide_layouts:
        if layout.name in layout_name_candidates:
            return layout
    return None


def get_layout(prs, layout_type):
    """レイアウトタイプに応じたスライドレイアウトを取得する。

    燈テンプレート（Google Slides由来）のレイアウト名に対応。
    """
    layouts = list(prs.slide_layouts)

    if layout_type == "title":
        found = find_layout(prs, [
            "TITLE", "TITLE_2", "TITLE_1",
            "Title Slide", "タイトル スライド",
        ])
        return found or layouts[0] if layouts else None

    elif layout_type == "section":
        found = find_layout(prs, [
            "SECTION_HEADER",
            "Section Header", "セクション見出し",
        ])
        return found or (layouts[3] if len(layouts) > 3 else layouts[0])

    elif layout_type == "two_column":
        found = find_layout(prs, [
            "TITLE_AND_TWO_COLUMNS",
            "Title and Two Columns",
        ])
        return found or (layouts[6] if len(layouts) > 6 else layouts[0])

    elif layout_type in ("content", "table", "closing"):
        found = find_layout(prs, [
            "TITLE_AND_BODY",
            "Title and Content", "タイトルとコンテンツ",
        ])
        return found or (layouts[4] if len(layouts) > 4 else layouts[0])

    elif layout_type == "blank":
        found = find_layout(prs, ["BLANK", "Blank", "白紙"])
        return found or layouts[-1]

    # フォールバック
    return layouts[4] if len(layouts) > 4 else layouts[0]


def add_title_slide(prs, slide_data):
    """タイトルスライドを追加"""
    layout = get_layout(prs, "title")
    slide = prs.slides.add_slide(layout)

    # タイトル
    if slide.shapes.title:
        set_text_with_font(
            slide.shapes.title.text_frame,
            slide_data.get("title", ""),
            font_size=Pt(28),
            bold=True,
        )

    # サブタイトル（2番目のプレースホルダー）
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:  # タイトル以外
                set_text_with_font(
                    ph.text_frame, subtitle,
                    font_size=Pt(16),
                    color=COLOR_LIGHT,
                )
                break

    return slide


def add_section_slide(prs, slide_data):
    """セクション区切りスライドを追加"""
    layout = get_layout(prs, "section")
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        set_text_with_font(
            slide.shapes.title.text_frame,
            slide_data.get("title", ""),
            font_size=Pt(24),
            bold=True,
        )

    return slide


def add_content_slide(prs, slide_data):
    """コンテンツスライドを追加（タイトル + Markdown本文）"""
    layout = get_layout(prs, "content")
    slide = prs.slides.add_slide(layout)

    # タイトル
    if slide.shapes.title:
        set_text_with_font(
            slide.shapes.title.text_frame,
            slide_data.get("title", ""),
            font_size=Pt(22),
            bold=True,
        )

    # 本文
    body = slide_data.get("body", "")
    if body:
        # コンテンツ用プレースホルダーを探す
        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                body_ph = ph
                break

        if body_ph:
            paragraphs = parse_markdown_to_paragraphs(body)
            apply_paragraphs_to_textframe(body_ph.text_frame, paragraphs)
        else:
            # プレースホルダーがない場合はテキストボックスを追加
            _add_body_textbox(slide, body)

    return slide


def add_two_column_slide(prs, slide_data):
    """2カラムスライドを追加（TITLE_AND_TWO_COLUMNS レイアウト使用）"""
    layout = get_layout(prs, "two_column")
    slide = prs.slides.add_slide(layout)

    # タイトル (idx=0)
    if slide.shapes.title:
        set_text_with_font(
            slide.shapes.title.text_frame,
            slide_data.get("title", ""),
            font_size=Pt(22),
            bold=True,
        )

    # プレースホルダーをインデックスで取得
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # 左カラム (idx=1, BODY)
    left_text = slide_data.get("left", "")
    if left_text and 1 in phs:
        paragraphs = parse_markdown_to_paragraphs(left_text)
        apply_paragraphs_to_textframe(phs[1].text_frame, paragraphs)

    # 右カラム (idx=2, BODY)
    right_text = slide_data.get("right", "")
    if right_text and 2 in phs:
        paragraphs = parse_markdown_to_paragraphs(right_text)
        apply_paragraphs_to_textframe(phs[2].text_frame, paragraphs)

    # サブタイトル (idx=3, 4) - 左右のカラムタイトル
    left_title = slide_data.get("left_title", "")
    right_title = slide_data.get("right_title", "")
    if left_title and 3 in phs:
        set_text_with_font(phs[3].text_frame, left_title, font_size=Pt(12), bold=True)
    if right_title and 4 in phs:
        set_text_with_font(phs[4].text_frame, right_title, font_size=Pt(12), bold=True)

    return slide


def add_table_slide(prs, slide_data):
    """テーブルスライドを追加"""
    layout = get_layout(prs, "content")
    slide = prs.slides.add_slide(layout)

    # タイトル
    if slide.shapes.title:
        set_text_with_font(
            slide.shapes.title.text_frame,
            slide_data.get("title", ""),
            font_size=Pt(22),
            bold=True,
        )

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])

    if not headers:
        return slide

    # 既存のコンテンツプレースホルダーをクリア
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            ph.text_frame.clear()

    num_rows = len(rows) + 1  # ヘッダー含む
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.8),
        Inches(9.0), Inches(0.4 + 0.35 * num_rows)
    )
    table = table_shape.table

    # ヘッダー
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ACCENT

    # データ行
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            if col_idx >= num_cols:
                break
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(10)
                    run.font.color.rgb = COLOR_PRIMARY

    return slide


def add_closing_slide(prs, slide_data):
    """クロージングスライドを追加"""
    layout = get_layout(prs, "section")
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        set_text_with_font(
            slide.shapes.title.text_frame,
            slide_data.get("title", "ありがとうございました"),
            font_size=Pt(28),
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    body = slide_data.get("body", "")
    if body:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                set_text_with_font(
                    ph.text_frame, body,
                    font_size=Pt(14),
                    color=COLOR_LIGHT,
                    alignment=PP_ALIGN.CENTER,
                )
                break

    return slide


def _add_body_textbox(slide, body_text):
    """スライドにMarkdown本文のテキストボックスを追加する（フォールバック用）。"""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5)
    )
    txBox.text_frame.word_wrap = True
    paragraphs = parse_markdown_to_paragraphs(body_text)
    apply_paragraphs_to_textframe(txBox.text_frame, paragraphs)


# スライドタイプ → 関数のマッピング
SLIDE_BUILDERS = {
    "title": add_title_slide,
    "section": add_section_slide,
    "content": add_content_slide,
    "two_column": add_two_column_slide,
    "table": add_table_slide,
    "closing": add_closing_slide,
}


def generate(proposal_data_path):
    """proposal_data.json を読み込んでPPTXを生成する。"""
    with open(proposal_data_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    # テンプレート読み込み
    template_path = data.get("template", "")
    if template_path:
        full_template_path = PROJECT_ROOT / template_path
        if full_template_path.exists():
            prs = Presentation(str(full_template_path))
            print(f"テンプレート読み込み: {full_template_path}")
        else:
            print(f"警告: テンプレート {full_template_path} が見つかりません。ブランクで生成します。")
            prs = Presentation()
    else:
        prs = Presentation()

    # テンプレートのスライドサイズをそのまま使用（上書きしない）
    # 燈テンプレートは 10.0 x 5.6 inches (Google Slides由来)

    # テンプレートの既存スライドを削除（レイアウトは保持される）
    existing_count = len(prs.slides)
    if existing_count > 0:
        print(f"テンプレートの既存スライド {existing_count} 枚を削除...")
        for i in range(existing_count - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

    # レイアウト一覧を表示（デバッグ用）
    print(f"利用可能なレイアウト:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")

    # スライド生成
    slides = data.get("slides", [])
    for i, slide_data in enumerate(slides):
        layout_type = slide_data.get("layout", "content")
        builder = SLIDE_BUILDERS.get(layout_type, add_content_slide)
        builder(prs, slide_data)
        print(f"  スライド {i+1}: {layout_type} - {slide_data.get('title', '(no title)')}")

    # 出力
    output_path = PROJECT_ROOT / data.get("output", "output.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nPPTX生成完了: {output_path}")
    return str(output_path)


def main():
    if len(sys.argv) < 2:
        print("使い方: python scripts/generate_proposal.py <proposal_data.json>")
        sys.exit(1)

    proposal_data_path = sys.argv[1]
    if not os.path.exists(proposal_data_path):
        print(f"エラー: {proposal_data_path} が見つかりません")
        sys.exit(1)

    generate(proposal_data_path)


if __name__ == "__main__":
    main()
